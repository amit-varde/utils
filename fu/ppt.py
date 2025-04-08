import sys
import json
import os
import argparse
import base64
import csv
from io import StringIO, BytesIO
from datetime import datetime
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("Error: The 'python-pptx' module is not installed.")
    print("Please install it using one of the following commands:")
    print("  pip install python-pptx")
    print("  conda install -c conda-forge python-pptx")
    sys.exit(1)

class PPTXSlides:
    def __init__(self, file_path=None):
        """
        Initialize the PPTXSlides with empty data structures and optionally load a presentation.
        
        Args:
            file_path: Optional path to a PowerPoint file to load immediately
        """
        print("Executing: __init__")
        self.presentation = None  # The presentation object
        self.file_path = None  # Path to the loaded file
        self.slides = []  # All slides in presentation
        self.slide_layouts = []  # All slide layouts in the presentation
        self.title_slides = []  # Indices of title slides
        self.summary_slides = []  # Indices of summary slides
        self.section_header = []  # Indices of section header slides
        # Interesting Data Objects Stored in the PPTX
        self.section_summaries = {}  # Dictionary mapping section header indices to title and summary text
        self.summary_tables = {}  # Dictionary mapping slide titles to tables
        # If file_path is provided, load the presentation
        if file_path:
            self.load_presentation(file_path)
            
    def extract_presentation_metadata(self):
        """
        Extract all presentation metadata including summary tables.
        This is a convenience method that calls all extract methods.
        """
        # Extract summary tables from identified summary slides
        if self.summary_slides:
            self.extract_summary_slides_table()
            
        # Extract section summary information
        if self.section_header:
            self.extract_section_summary_slides()
        
        return {
            "slides": len(self.slides),
            "layouts": len(self.slide_layouts),
            "title_slides": self.title_slides,
            "summary_slides": self.summary_slides,
            "section_header_slides": self.section_header
        }
        
    def extract_summary_slides_table(self):
        """
        Extract table data from all identified summary slides.
        Populates the summary_tables dictionary with table data.
        """
        print("Executing: extract_summary_slides_table")
        for slide_idx in self.summary_slides:
            tables = self.extract_table_from_slide(slide_idx)
            if tables:
                self.summary_tables[slide_idx] = tables

    def extract_section_summary_slides(self):
        """
        Extract title and summary text from all identified section header slides.
        Populates the section_summaries dictionary.
        """
        print("Executing: extract_section_summary_slides")
        for section_idx in self.section_header:
            slide = self.slides[section_idx]
            section_title = ""
            section_summary = ""
            
            for shape_idx, shape in enumerate(slide.shapes):
                if shape_idx == 0 and hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                    section_title = shape.text_frame.text
                elif shape_idx == 1 and hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                    section_summary = shape.text_frame.text
            
            # Store in section_summaries dictionary
            self.section_summaries[section_idx] = {
                "section_title": section_title,
                "section_summary": section_summary
            }
            print(f"Extracted section info from slide {section_idx + 1}: '{section_title}'")

    # ------------------------------------------------
    # File handling methods
    # ------------------------------------------------
    
    def load_presentation(self, file_path):
        """
        Load a PowerPoint presentation from a file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of slide objects
        """
        print("Executing: load_presentation")
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {file_path}")
        if not os.path.isfile(file_path):
            raise ValueError(f"The path is not a file: {file_path}")
        print(f"Attempting to process file: {file_path}")
        print(f"File exists: {os.path.exists(file_path)}")
        print(f"File size: {os.path.getsize(file_path)} bytes")
        # Use python-pptx to read the presentation
        try:
            self.file_path = file_path
            self.presentation = Presentation(file_path)
            self.slides = list(self.presentation.slides)
            self.slide_layouts = list(self.presentation.slide_layouts)
            # Identify title slides (first slide and slides with title layout)
            self.title_slides = [0]  # First slide is typically a title slide
            for i, slide in enumerate(self.slides):
                if i > 0 and hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
                    if 'title' in slide.slide_layout.name.lower():
                        self.title_slides.append(i)
                        
            # Identify section header slides
            self.section_header = self._identify_section_header_slides()
            
            # Identify summary slides
            self.summary_slides = self._identify_summary_slides()
            return self.slides
        except Exception as e:
            print(f"Error reading PowerPoint file: {e}")
            raise
    
    def save_presentation(self, output_file):
        """
        Save a PowerPoint presentation to a file.
        
        Args:
            output_file: Path to save the PowerPoint file
        
        Returns:
            Path to the saved file
        """
        print("Executing: save_presentation")
        output_dir = os.path.dirname(output_file)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir)
        print(f"Writing presentation to: {output_file}")
        self.presentation.save(output_file)
        print(f"Successfully saved presentation to: {output_file}")
        return output_file

    # ------------------------------------------------
    # Private identification methods
    # ------------------------------------------------
    
    def _identify_section_header_slides(self):
        """
        Identify section header slides based on their layout name.
        
        Returns:
            List of slide indices that are section headers
        """
        print("Executing: _identify_section_header_slides")
        section_headers = []
        
        # Iterate through slides to find those with Section Header layout
        for i, slide in enumerate(self.slides):
            if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
                if slide.slide_layout.name == "Section Header":
                    section_headers.append(i)
                    print(f"Found section header slide at index {i}: {self.extract_metadata_from_slide(slide, i)['title']}")
                    
                    # Extract section title and summary
                    section_title = ""
                    section_summary = ""
                    
                    for shape_idx, shape in enumerate(slide.shapes):
                        if shape_idx == 0 and hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                            section_title = shape.text_frame.text
                        elif shape_idx == 1 and hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                            section_summary = shape.text_frame.text
                    
                    # Store in section_summaries dictionary
                    self.section_summaries[i] = {
                        "section_title": section_title,
                        "section_summary": section_summary
                    }
        
        return section_headers
    
    def _identify_summary_slides(self):
        """
        Identify summary slides where the title starts with "Summary" (case insensitive)
        and the slide contains at least one table.
        
        Returns:
            List of slide indices that are summary slides
        """
        print("Executing: _identify_summary_slides")
        summary_slides = []
        for i, slide in enumerate(self.slides):
            is_summary_slide = False
            has_table = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    if hasattr(shape, "is_title") and shape.is_title and shape.text.lower().startswith("summary"):
                        is_summary_slide = True
                        break
                    elif shape.text.lower().startswith("summary"):
                        is_summary_slide = True
                        break
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    has_table = True
                    break
            if is_summary_slide and has_table:
                summary_slides.append(i)
        return summary_slides
    
    def _extract_color(self, color_obj):
        """Helper method to extract color information."""
        print("Executing: _extract_color")
        if not color_obj:
            return None
        color_dict = {}
        if hasattr(color_obj, 'rgb'):
            rgb = color_obj.rgb
            if rgb:
                color_dict['rgb'] = f'#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}'
        if hasattr(color_obj, 'theme_color'):
            color_dict['theme_color'] = str(color_obj.theme_color)
        return color_dict

    # ------------------------------------------------
    # Extract methods
    # ------------------------------------------------
    def extract_metadata_from_slide(self, slide, index):
        """
        Extract metadata from a single slide.
        
        Args:
            slide: PowerPoint slide object
            index: The slide index (0-based)
            
        Returns:
            Dictionary containing slide metadata
        """
        print(f"Extracting metadata for slide {index + 1}")
        slide_number = index + 1
        layout_name = "Unknown"
        if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
            layout_name = slide.slide_layout.name
        slide_title = "Untitled"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if hasattr(shape, "is_title") and shape.is_title:
                    slide_title = shape.text
                    break
                elif slide_title == "Untitled":
                    slide_title = shape.text
        shape_counts = {}
        for shape in slide.shapes:
            shape_type = str(shape.shape_type)
            if shape_type in shape_counts:
                shape_counts[shape_type] += 1
            else:
                shape_counts[shape_type] = 1
        metadata = {
            "slide_number": slide_number,
            "layout_name": layout_name,
            "title": slide_title,
            "shape_counts": shape_counts,
            "total_shapes": len(slide.shapes)
        }
        return metadata
    
    def extract_table_from_slide(self, index):
        """
        Extract table data from a specific slide.
        
        Args:
            index: The index of the slide (0-based)
            
        Returns:
            List of tables found on the slide
        """
        print(f"Executing: extract_table_from_slide for slide {index + 1}")
        
        if index >= len(self.slides):
            print(f"Invalid slide index: {index}, max index is {len(self.slides)-1}")
            return []
            
        slide = self.slides[index]
        slide_tables = []
        slide_title = self.extract_metadata_from_slide(slide, index)["title"]
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                table_data = []
                for r in range(len(table.rows)):
                    row_data = []
                    for c in range(len(table.columns)):
                        cell = table.cell(r, c)
                        cell_text = cell.text.replace('\n', ' ')
                        row_data.append(cell_text)
                    table_data.append(row_data)
                slide_tables.append(table_data)
        
        # Store in summary_tables if there are tables found
        if slide_tables:
            self.summary_tables[index] = slide_tables
        
        return slide_tables
    
    # ------------------------------------------------
    # Get methods
    # ------------------------------------------------
    
    def get_summary_table(self, index):
        """
        Get summary table data by slide index.
        
        Args:
            index: Index of the slide containing the table (0-based)
            
        Returns:
            Tables data for the specified slide, or None if no tables found
        """
        # Validate the index
        if not (0 <= index < len(self.slides)):
            print(f"Warning: Invalid slide index {index}, valid range is 0-{len(self.slides)-1}")
            return None
        
        # Extract the slide title from the index
        slide_title = self.extract_metadata_from_slide(self.slides[index], index)["title"]
        
        # Extract tables if they don't exist for this slide
        if index not in self.summary_tables:
            self.extract_table_from_slide(index)
        
        # Return tables for the slide index
        if index in self.summary_tables:
            return self.summary_tables[index]
        else:
            print(f"Warning: No tables found for slide at index {index}")
            return None
    
    def get_section_summary(self, index):
        """
        Get section summary information by section index.
        
        Args:
            index: Index of the section header slide
            
        Returns:
            Dictionary with section_title and section_summary or None if not found
        """
        if index in self.section_summaries:
            return self.section_summaries[index]
        else:
            print(f"Warning: No section summary found for index {index}")
            return None

    # ------------------------------------------------
    # Set methods
    # ------------------------------------------------
    
    def set_presentation_title_slide(self, index=0, title_text=None, attribution_text=None):
        """
        Updates a slide with title and attribution text.
        
        Args:
            index: Index of the slide to update (0-based, default is first slide)
            title_text: Text to use for the title (if None, uses "Weekly Update" with today's date)
            attribution_text: Text to use for attribution (if None, uses default attribution)
            
        Returns:
            True if title was updated, False otherwise
        """
        print(f"Executing: set_presentation_title_slide for slide {index + 1}")
        if not self.presentation or not self.slides or not (0 <= index < len(self.slides)):
            print(f"Invalid slide index: {index}")
            return False
            
        slide = self.slides[index]
        text_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_shapes.append(shape)
                
        today = datetime.now().strftime("%m-%d-%Y")
        if title_text is None:
            title_text = f"Weekly Update {today}"
        if attribution_text is None:
            attribution_text = "Alec .. did this automatically"
            
        content = [title_text, attribution_text]
        for i, text in enumerate(content):
            if text is None:
                continue
            if i < len(text_shapes):
                text_shapes[i].text_frame.text = text
            elif i == 1 and len(text_shapes) == 1:
                p = text_shapes[0].text_frame.add_paragraph()
                p.text = f"\n\n{attribution_text}"
            else:
                print(f"Warning: Not enough text shapes on slide {index + 1} for content #{i+1}")
        return True
    
    def set_summary_table(self, index, row, col, new_text):
        """
        Update a specific cell in a table on a slide.
        
        Args:
            index: Index of the slide containing the table (0-based)
            row: Row index (0-based)
            col: Column index (0-based)
            new_text: New text to set in the cell
        
        Returns:
            Boolean indicating success or failure
        """
        # Validate the index
        if not (0 <= index < len(self.slides)):
            print(f"Invalid slide index: {index}, valid range is 0-{len(self.slides)-1}")
            return False
        
        # Extract tables if they don't exist for this slide
        if index not in self.summary_tables:
            tables = self.extract_table_from_slide(index)
            if not tables:
                print(f"No tables found on slide at index {index}")
                return False
        
        # Get the tables for this slide
        slide_tables = self.summary_tables[index]
        if not slide_tables:
            print(f"No tables found on slide at index {index}")
            return False
            
        table_data = slide_tables[0]  # Use the first table
        if row < 0 or row >= len(table_data):
            print(f"Invalid row index: {row}, valid range is 0-{len(table_data)-1}")
            return False
        if col < 0 or col >= len(table_data[0]):
            print(f"Invalid column index: {col}, valid range is 0-{len(table_data[0])-1}")
            return False
            
        # Update the data in memory
        table_data[row][col] = new_text
        
        # Update the actual PowerPoint object
        slide = self.slides[index]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    shape.table.cell(row, col).text = new_text
                    return True
                except IndexError:
                    print(f"Error updating PowerPoint table: Index out of range")
                    return False
                    
        print(f"Updated table data in memory but couldn't update PowerPoint object")
        return True
    
    def set_section_summary(self, index, title=None, summary=None):
        """
        Update section title and/or summary text.
        
        Args:
            index: Index of the section header slide
            title: New title text (optional)
            summary: New summary text (optional)
            
        Returns:
            Boolean indicating success or failure
        """
        # This method already matches the requested signature, but we'll ensure it works properly
        if index not in self.section_summaries:
            print(f"Error: No section found with index {index}")
            return False
            
        # Update in-memory data
        if title is not None:
            self.section_summaries[index]["section_title"] = title
            
        if summary is not None:
            self.section_summaries[index]["section_summary"] = summary
            
        # Update actual PowerPoint slide
        if index < len(self.slides):
            slide = self.slides[index]
            for shape_idx, shape in enumerate(slide.shapes):
                if shape_idx == 0 and hasattr(shape, "text_frame") and title is not None:
                    shape.text_frame.text = title
                elif shape_idx == 1 and hasattr(shape, "text_frame") and summary is not None:
                    shape.text_frame.text = summary
            return True
        else:
            print(f"Warning: Could update section summary in memory but not in PowerPoint (invalid slide index)")
            return False

    # ------------------------------------------------
    # Show/display methods
    # ------------------------------------------------
    
    def show_slide_metadata(self, slide, index):
        """
        Display metadata for a specific slide using extract_metadata_from_slide.
        
        Args:
            slide: PowerPoint slide object
            index: The slide index (0-based)
        """
        metadata = self.extract_metadata_from_slide(slide, index)
        print(f"Slide #{metadata['slide_number']}:")
        print(f"  Layout: {metadata['layout_name']}")
        print(f"  Title: {metadata['title']}")
        print(f"  Total shapes: {metadata['total_shapes']}")
        print("  Shape counts:")
        for shape_type, count in metadata['shape_counts'].items():
            print(f"    {shape_type}: {count}")
        # Display names/texts of shapes on the slide
        print("  Shapes:")
        for i, shape in enumerate(slide.shapes):
            # Try to extract shape name or identifier
            shape_name = "Unnamed"
            if hasattr(shape, "name") and shape.name:
                shape_name = shape.name
            # Try to extract text from the shape
            shape_text = ""
            if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                shape_text = shape.text_frame.text
                if len(shape_text) > 30:
                    shape_text = shape_text[:27] + "..."
                if shape_text:
                    shape_text = f" - \"{shape_text}\""
            # Get shape type name
            shape_type = str(shape.shape_type)
            print(f"    Shape #{i+1}: {shape_name} ({shape_type}){shape_text}")
        # Check if this is a summary slide and display tables if it is
        if index in self.summary_slides:
            print("  This is a summary slide with tables.")
            # Make sure summary tables are extracted first
            if not self.summary_tables:
                self.extract_table_from_slide(index)
            slide_title = metadata['title']
            self.show_summary_tables(slide_title)
        print("-" * 50)
    
    def show_presentation_metadata(self):
        """
        Display metadata for all slides in the presentation by
        looping through each slide and running show_slide_metadata.
        """
        print("Executing: show_presentation_metadata")
        print(f"\nPresentation contains {len(self.slides)} slides")
        print("=" * 50)
        for i, slide in enumerate(self.slides):
            self.show_slide_metadata(slide, i)
    
    def show_summary_tables(self, slide_title):
        """
        Display table data in a formatted text table for a specific slide title.
        
        Args:
            slide_title: Title of the slide whose tables to display
        """
        print("Executing: show_summary_tables")
        if slide_title not in self.summary_tables:
            print(f"No table data found for slide with title: {slide_title}")
            return
        slide_tables = self.summary_tables[slide_title]
        print(f"\nTable data from slide ({slide_title}):")
        for table_idx, table_data in enumerate(slide_tables):
            print(f"Table {table_idx + 1}:")
            try:
                if not table_data:
                    print("  Empty table")
                    continue
                col_widths = [0] * len(table_data[0])
                for row in table_data:
                    for i, cell in enumerate(row):
                        if i < len(col_widths):
                            col_widths[i] = max(col_widths[i], len(cell))
                separator = "  +" + "+".join("-" * (width + 2) for width in col_widths) + "+"
                print(separator)
                for row in table_data:
                    row_text = "  |"
                    for i, cell in enumerate(row):
                        if i < len(col_widths):
                            padding = col_widths[i] - len(cell)
                            row_text += f" {cell}{' ' * padding} |"
                    print(row_text)
                    print(separator)
            except Exception as e:
                print(f"Error formatting table: {e}")
                print(table_data)
            print()

def main():
    print("Executing: main")
    parser = argparse.ArgumentParser(description='Extract information from PowerPoint files.')
    parser.add_argument('--file', '-f', type=str, 
                       default="/Users/amit/work/utils/docprocessors/test_cases/3-27-2025.pptx",
                       help='Path to the PowerPoint file')
    parser.add_argument('--output', '-o', type=str,
                       help='Output file path (optional, default: input_update.pptx)')
    args = parser.parse_args()
    file_path = args.file
    output_path = args.output
    
    try:
        if not os.path.exists(file_path):
            print(f"Error: File not found: {file_path}")
            sys.exit(1)
        print(f"Reading presentation from: {file_path}")
        status_pptx = PPTXSlides(file_path)
        status_pptx.extract_presentation_metadata()
        status_pptx.show_presentation_metadata()
        
        # Update the title slide
        today = datetime.now().strftime("%m-%d-%Y")
        status_pptx.set_presentation_title_slide(title_text=f"{today}: Weekly Update")
        
        # Save the presentation
        if not output_path:
            base_name, ext = os.path.splitext(file_path)
            output_path = f"{base_name}_update{ext}"
        status_pptx.save_presentation(output_path)
        print(f"PPTX file updated successfully from {file_path} to {output_path}")
    except Exception as e:
        print(f"Error processing presentation: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()
